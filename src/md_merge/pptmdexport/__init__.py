from __future__ import annotations

import argparse
import json
import logging
from pathlib import Path

from pptx import Presentation

from md_merge._filters import escape_backslash_smart, replace_slide_placeholders
from md_merge._note_parser import parse_note_md_blocks
from md_merge._output import EXIT_BAD_INPUT, EXIT_FAILURE, EXIT_NOT_FOUND, EXIT_OK, emit, setup_logging
from md_merge.merge._recipe import load_yaml, resolve_input_path, resolve_workdir, setup_recipe_file_logging
from md_merge.pptmerge._config import resolve_base


def _normalize_cropsrc(cropsrc: dict | list) -> list[dict]:
    """cropsrc を常にリストとして返す（辞書は単一要素リストに変換）。"""
    if isinstance(cropsrc, list):
        return cropsrc
    if isinstance(cropsrc, dict):
        return [cropsrc]
    return []


def _resolve_pptx_path(val: str, yaml_path: Path, workdir: Path | None) -> Path:
    base = resolve_base(yaml_path, workdir)
    p = Path(val)
    return p if p.is_absolute() else (base / p).resolve()


def run(args: argparse.Namespace) -> int:
    setup_logging(args.log_level)

    yaml_path = resolve_input_path(args)
    if yaml_path is None:
        logging.error("No input file specified")
        return EXIT_BAD_INPUT
    if not yaml_path.exists():
        logging.error("Input file not found: %s", yaml_path)
        return EXIT_NOT_FOUND

    recipe = load_yaml(yaml_path)
    raw_cropsrc = (recipe.get("input") or {}).get("cropsrc") or {}
    entries = _normalize_cropsrc(raw_cropsrc)
    output_section = recipe.get("output") or {}
    cli_workdir = Path(args.workdir).resolve() if getattr(args, "workdir", None) else None
    workdir = resolve_workdir(recipe, yaml_path, cli_workdir)
    setup_recipe_file_logging(recipe, yaml_path, workdir)
    recipe_force = bool(output_section.get("force"))
    force = getattr(args, "force", False) or recipe_force
    dry_run = getattr(args, "dry_run", False)
    json_mode = getattr(args, "json", False)

    if not entries:
        logging.error("input.cropsrc が指定されていません: %s", yaml_path)
        return EXIT_BAD_INPUT

    # ── 全エントリのパス解決と存在確認（fail fast）────────────────────────
    resolved_pptx_paths: list[Path] = []
    for i, entry in enumerate(entries):
        label = f"input.cropsrc[{i}]" if len(entries) > 1 else "input.cropsrc"
        pptx_val = entry.get("pptx")

        if not pptx_val:
            logging.error("%s.pptx が指定されていません: %s", label, yaml_path)
            return EXIT_BAD_INPUT

        pptx_path = _resolve_pptx_path(str(pptx_val), yaml_path, workdir)
        if not pptx_path.exists():
            logging.error("PPTX ファイルが見つかりません: %s", pptx_path)
            return EXIT_NOT_FOUND

        resolved_pptx_paths.append(pptx_path)

    # ── エントリごとに処理 ────────────────────────────────────────────────
    errors = 0
    json_entries: list[dict] = []

    for pptx_path in resolved_pptx_paths:
        try:
            prs = Presentation(str(pptx_path))
        except Exception as e:
            logging.error("PPTX ファイルを開けません: %s: %s", pptx_path, e)
            errors += 1
            continue

        tasks: list[dict] = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            note_text = ""
            if slide.has_notes_slide:
                tf = slide.notes_slide.notes_text_frame
                note_text = (tf.text or "") if tf else ""

            title_shape = slide.shapes.title
            slide_title = title_shape.text if title_shape is not None else ""
            note_text = replace_slide_placeholders(note_text, pptx_path, slide_title, slide_idx)

            for body in parse_note_md_blocks(note_text):
                lines = body.splitlines()
                if not lines:
                    continue

                # 1行目: export-note 指示
                first_line = lines[0]
                if not first_line.startswith("export-note:"):
                    logging.warning(
                        "スライド %d: export-note が1行目にありません: %s", slide_idx, first_line
                    )
                    continue

                _, _, export_note_val = first_line.partition(":")
                export_note_val = export_note_val.strip()

                if not export_note_val.lower().endswith(".md"):
                    logging.warning(
                        "スライド %d: export-note のファイル名が .md ではありません: %s",
                        slide_idx, export_note_val,
                    )
                    continue

                # 2行目以降: MD テキスト
                md_text = "\n".join(lines[1:])
                display_line = lines[1] if len(lines) > 1 else ""

                export_path = Path(export_note_val)
                out_path = export_path if export_path.is_absolute() else (pptx_path.parent / export_path).resolve()

                tasks.append({
                    "slide": slide_idx,
                    "output": out_path,
                    "md_text": md_text,
                    "display_line": display_line,
                })

        if not tasks:
            if not json_mode:
                print(f"export-note を持つスライドが見つかりません: {pptx_path}")
            json_entries.append({"pptx": str(pptx_path), "tasks": [], "writes": []})
            continue

        if not json_mode:
            print(f"対象スライド ({len(tasks)} 件): {pptx_path}")
            for task in tasks:
                print(f"  スライド {task['slide']:>4}: {task['output']}  {task['display_line']}")

        write_results: list[dict] = []

        for task in tasks:
            slide_n: int = task["slide"]
            out_path: Path = task["output"]

            if out_path.exists() and not force:
                logging.error("出力ファイルが既に存在します (--force で上書き): %s", out_path)
                errors += 1
                continue

            if dry_run:
                emit(args, "ok", "pptmdexport", pptx_path, out_path, dry_run=True)
                write_results.append({"slide": slide_n, "output": str(out_path), "display_line": task["display_line"]})
                continue

            try:
                out_path.parent.mkdir(parents=True, exist_ok=True)
                with out_path.open("w", encoding="utf-8", newline="\n") as f:
                    f.write(escape_backslash_smart(task["md_text"]))
                emit(args, "ok", "pptmdexport", pptx_path, out_path)
                write_results.append({"slide": slide_n, "output": str(out_path), "display_line": task["display_line"]})
            except Exception as e:
                logging.error("書き込み失敗: slide %d: %s", slide_n, e)
                errors += 1

        if getattr(args, "closepptx", False):
            del prs

        json_entries.append({
            "pptx": str(pptx_path),
            "tasks": [{"slide": t["slide"], "output": str(t["output"]), "display_line": t["display_line"]} for t in tasks],
            "writes": write_results,
        })

    if json_mode:
        print(json.dumps({
            "status": "ok" if errors == 0 else "error",
            "command": "pptmdexport",
            "entries": json_entries,
        }, ensure_ascii=False, indent=2))

    return EXIT_OK if errors == 0 else EXIT_FAILURE
