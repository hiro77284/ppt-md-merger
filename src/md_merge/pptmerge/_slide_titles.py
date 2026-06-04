from __future__ import annotations

import logging
import re
import os
import win32com.client
from collections.abc import Iterator
from pathlib import Path
from typing import Any, TextIO

import pptx.presentation as pptx_presentation
import yaml
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from ._config import ConfigError
from ._slide_operation import SlideOperation, format_slide_operations_block


chapter_title_format = r'#CHAPTNUM##SEPA##CONTENT#'
section_title_format = r'#CHAPTNUM##DELM##SECTNUM##SEPA##CONTENT#'

NOREPLACE_MARKER = "#NOREPLACE#"
NONUMBERING_MARKER = "#NONUMBERING#"


def _apply_index_counter(current: int, value: Any) -> int:
    if value is None:
        return current + 1
    if isinstance(value, int) and not isinstance(value, bool):
        return int(value)
    if isinstance(value, float) and value.is_integer():
        return int(value)

    s = str(value).strip()
    if not s:
        return current + 1
    if s.isdigit():
        return int(s)
    if s.startswith("+"):
        rest = s[1:]
        if rest.isdigit():
            return current + int(rest)
        return current + 1 if rest == "" else current
    if s.startswith("-"):
        rest = s[1:]
        if rest.isdigit():
            return current - int(rest)
        return current - 1 if rest == "" else current
    return current


_INDEXER_DEFAULTS: dict[str, Any] = {
    "chapter_marker": "#CHAPT#",
    "section_marker": "#SECTION#",
    "stay_marker": "#STAY#",
    "delimiter": ".",
    "separator": ") ",
    "deletecsl": False,
    "deletecsp": False,
}


def _get_indexer_value(config: dict[str, Any] | None, key: str) -> Any:
    default = _INDEXER_DEFAULTS[key]
    if config is None:
        return default
    raw = config.get("indexer")
    if not isinstance(raw, dict):
        return default
    if key not in raw:
        return default
    val = raw[key]
    if val == "":
        return default
    return val


def format_indexer_config_block(config: dict[str, Any]) -> str:
    lines: list[str] = ["indexer 設定:"]
    raw = config.get("indexer")
    if not isinstance(raw, dict):
        lines.append("  （indexer セクションなし。参照用の既定値）")
        for key, default in _INDEXER_DEFAULTS.items():
            lines.append(f"  {key}: {default!r}")
        return "\n".join(lines) + "\n"
    for key, default in _INDEXER_DEFAULTS.items():
        if key in raw:
            lines.append(f"  {key}: {raw[key]!r}")
        else:
            lines.append(f"  {key}: {default!r} （既定）")
    return "\n".join(lines) + "\n"


def load_slide_titles(prs: pptx_presentation.Presentation) -> list[tuple[int, str]]:
    result: list[tuple[int, str]] = []
    for index, slide in enumerate(prs.slides, start=1):
        title_text = ""
        if slide.shapes.title is not None:
            title_text = (slide.shapes.title.text or "").strip()
        result.append((index, title_text))
    return result


def replace_slide_title(
    prs: pptx_presentation.Presentation,
    slide_number: int,
    title: str,
) -> None:
    n = len(prs.slides)
    if slide_number < 1 or slide_number > n:
        raise ValueError(f"スライド番号が範囲外です: {slide_number} （プレゼンは {n} 枚）")
    slide = prs.slides[slide_number - 1]
    if slide.shapes.title is None:
        raise ValueError(f"スライド {slide_number} にタイトルプレースホルダがありません。")

    tf = slide.shapes.title.text_frame
    txBody = tf._txBody
    paras = txBody.findall(qn("a:p"))

    if not paras:
        slide.shapes.title.text = title
        return

    # Keep only the first paragraph (preserves <a:pPr>: alignment, spacing, etc.)
    for p in paras[1:]:
        txBody.remove(p)

    para = paras[0]
    # Remove line breaks and field elements in this paragraph
    for br in para.findall(qn("a:br")):
        para.remove(br)
    for fld in para.findall(qn("a:fld")):
        para.remove(fld)

    runs = para.findall(qn("a:r"))
    if runs:
        # Set text on the first run (preserves <a:rPr>: font, size, bold, color)
        t_elem = runs[0].find(qn("a:t"))
        if t_elem is None:
            t_elem = etree.SubElement(runs[0], qn("a:t"))
        t_elem.text = title
        for r in runs[1:]:
            para.remove(r)
    else:
        r_elem = etree.SubElement(para, qn("a:r"))
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = title


def _iter_shape_text_fragments(shape: Any) -> Iterator[str]:
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            yield from _iter_shape_text_fragments(child)
        return
    if getattr(shape, "has_table", False) and shape.table is not None:
        for row in shape.table.rows:
            for cell in row.cells:
                tf = cell.text_frame
                if tf is not None:
                    t = (tf.text or "").strip()
                    if t:
                        yield t
        return
    if getattr(shape, "has_text_frame", False) and shape.text_frame is not None:
        t = (shape.text_frame.text or "").strip()
        if t:
            yield t


def load_slide_all_shape_texts(prs: pptx_presentation.Presentation) -> list[str]:
    combined: list[str] = []
    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            parts.extend(_iter_shape_text_fragments(shape))
        combined.append("".join(parts))
    return combined


def filter_slides_excluding_text_marker(
    src_path: Path,
    slide_list: list[int] | None,
    marker: str = "#CSL#",
) -> list[int]:
    prs = Presentation(str(src_path))
    texts = load_slide_all_shape_texts(prs)
    n = len(texts)

    if slide_list is None:
        indices = list(range(1, n + 1))
    else:
        indices = slide_list
        for i in indices:
            if not isinstance(i, int) or i < 1 or i > n:
                raise ConfigError(
                    f"slides のスライド番号がファイルの範囲外です: {src_path} （指定 {i}、ファイルは {n} 枚）"
                )
    if not marker:
        return list(indices)
    return [i for i in indices if marker not in texts[i - 1]]


def postprocess_after_merge(
    pptx_path: Path,
    log_file: Path,
    duplicate_to: TextIO | None = None,
    config: dict[str, Any] | None = None,
    operations: list[SlideOperation] = [],
) -> None:
    chapter_count = 0
    section_count = 0
    slides_modified = False

    merged_prs = Presentation(str(pptx_path))
    rows = load_slide_titles(merged_prs)
    lines: list[str] = []
    slide_texts = load_slide_all_shape_texts(merged_prs)

    chapter_marker = str(_get_indexer_value(config, "chapter_marker"))
    section_marker = str(_get_indexer_value(config, "section_marker"))
    stay_marker = str(_get_indexer_value(config, "stay_marker"))

    chapsectfreeze_flag = False
    chapt_or_section_flag = "chapter"

    for operation in operations:
        has_chapt = False
        has_sect = False
        has_stay = False
        if operation.action == "insertslide":
            title_disp = operation.slidetitle if operation.slidetitle else "(タイトルなし)"
            slide_body = slide_texts[operation.composedslidenum - 1]
            has_nonumbering = NONUMBERING_MARKER in slide_body
            has_chapt = bool(chapter_marker) and chapter_marker in slide_body
            has_sect = bool(section_marker) and section_marker in slide_body
            has_stay = bool(stay_marker) and stay_marker in slide_body

            if has_chapt and not chapsectfreeze_flag and not has_stay:
                chapt_or_section_flag = "chapter"
                chapter_count = _apply_index_counter(chapter_count, "+")
                section_count = 0
            if has_sect and not chapsectfreeze_flag and not has_stay:
                chapt_or_section_flag = "section"
                section_count = _apply_index_counter(section_count, "1")

            if not has_chapt and not has_sect:
                if not has_stay and not chapsectfreeze_flag:
                    chapt_or_section_flag = "section"
                    section_count = _apply_index_counter(section_count, "+")

        elif operation.action == "blank":
            pass
        elif operation.action == "chapter":
            chapter_count = _apply_index_counter(chapter_count, operation.value)
            section_count = _apply_index_counter(section_count, "0")
        elif operation.action == "section":
            section_count = _apply_index_counter(section_count, operation.value)
        elif operation.action == "beginstay":
            chapsectfreeze_flag = True
        elif operation.action == "endstay":
            chapsectfreeze_flag = False

        if operation.action == "insertslide":
            title_disp = operation.slidetitle if operation.slidetitle else "(タイトルなし)"
            if has_chapt or chapt_or_section_flag == "chapter":
                newtitle = (
                    chapter_title_format
                    .replace("#CHAPTNUM#", str(chapter_count))
                    .replace("#SEPA#", str(_get_indexer_value(config, "separator")))
                    .replace("#CONTENT#", title_disp)
                )
            elif has_sect or chapt_or_section_flag == "section":
                newtitle = (
                    section_title_format
                    .replace("#CHAPTNUM#", str(chapter_count))
                    .replace("#DELM#", str(_get_indexer_value(config, "delimiter")))
                    .replace("#SECTNUM#", str(section_count))
                    .replace("#SEPA#", str(_get_indexer_value(config, "separator")))
                    .replace("#CONTENT#", title_disp)
                )
            else:
                newtitle = title_disp

            if not has_nonumbering:
                try:
                    replace_slide_title(merged_prs, operation.composedslidenum, newtitle)
                    slides_modified = True
                except ValueError:
                    pass


    ppt_saved_by_com_shape_delete = _run_modify_slide_shapes(pptx_path, merged_prs, config)

    for i, title in rows:
        lines.append(f"[{i}] {title}" if title else f"[{i}] (タイトルなし)")

    block = "出力スライド一覧:\n" + "\n".join(lines) + "\n"
    if duplicate_to is not None:
        print(block, end="", file=duplicate_to)

    log_file.parent.mkdir(parents=True, exist_ok=True)
    with log_file.open("a", encoding="utf-8") as f:
        f.write(block)

    if slides_modified:
        if not ppt_saved_by_com_shape_delete:
            merged_prs.save(str(pptx_path))


_CHAPT_SECTION_STAY_MARKERS = (
    re.compile(r"#CHAPT#\s?.*"),
    re.compile(r"#SECTION#\s?.*"),
    re.compile(r"#STAY#\s?.*"),
)


def _strip_chapter_section_stay_markers_from_text_frame(tf: Any) -> None:
    for paragraph in tf.paragraphs:
        full = paragraph.text
        if "#CHAPT#" not in full and "#SECTION#" not in full and "#STAY#" not in full:
            continue
        new_full = full
        for pat in _CHAPT_SECTION_STAY_MARKERS:
            new_full = pat.sub("", new_full)
        if new_full != full:
            paragraph.text = new_full


def _strip_chapter_section_stay_markers_from_shape(shape: Any) -> None:
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _strip_chapter_section_stay_markers_from_shape(child)
        return
    if getattr(shape, "has_text_frame", False) and shape.text_frame is not None:
        _strip_chapter_section_stay_markers_from_text_frame(shape.text_frame)


def replace_CSLandCSP(run: Any, deletecsl: bool, deletecsp: bool) -> tuple[bool, bool]:
    tobreak = False
    needshapedeletion = False
    if re.search(r'#CSP#', run.text):
        if deletecsp:
            needshapedeletion = True
            tobreak = True
        else:
            before = run.text
            run.text = re.sub(r'#CSP#\s?', '', run.text)
    if re.search(r'#TEMP#', run.text) or re.search(r'#MEMO#', run.text):
        needshapedeletion = True
        tobreak = True
    if re.search(r'#CSL#', run.text):
        if deletecsl:
            tobreak = True
        else:
            run.text = re.sub(r'#CSL#\s?.*', '', run.text)
    return tobreak, needshapedeletion


def _run_modify_slide_shapes(
    pptx_path: Path,
    prs: Any,
    config: dict[str, Any] | None,
) -> bool:
    """Run marker stripping / shape deletion (``modify_slide_shapes``) using *config* indexer flags."""
    deletecsl = _get_indexer_value(config, "deletecsl")
    deletecsp = _get_indexer_value(config, "deletecsp")
    return modify_slide_shapes(pptx_path, prs, deletecsl, deletecsp)


def _collect_com_shape_text_for_markers(shape: Any) -> str:
    """Return concatenated text from a leaf COM shape (text frame + table cells)."""
    parts: list[str] = []
    try:
        if bool(shape.HasTextFrame) or shape.HasTextFrame == -1:
            if shape.TextFrame.HasText:
                parts.append(shape.TextFrame.TextRange.Text or "")
    except Exception:
        pass
    try:
        if shape.HasTable:
            table = shape.Table
            for r in range(1, table.Rows.Count + 1):
                for c in range(1, table.Columns.Count + 1):
                    try:
                        parts.append(table.Cell(r, c).Shape.TextFrame.TextRange.Text or "")
                    except Exception:
                        pass
    except Exception:
        pass
    return "".join(parts)


def _delete_com_shapes_with_markers(shapes: Any, deletecsp: bool) -> int:
    """Recursively delete shapes that contain markers. Returns count deleted.

    Groups: if the group name contains #CSP# (and deletecsp), the entire group
    is deleted without recursing into children. Otherwise children are checked.
    Leaves: text content is checked for #CSP#/#TEMP#/#MEMO#.
    """
    deleted = 0
    for shape in reversed(list(shapes)):
        try:
            if shape.Type == 6:  # msoGroup
                if deletecsp and "#CSP#" in (shape.Name or ""):
                    logging.debug(
                        "win32com: #CSP# 検出 → グループ削除 (name=%r)", shape.Name
                    )
                    shape.Delete()
                    deleted += 1
                else:
                    deleted += _delete_com_shapes_with_markers(shape.GroupItems, deletecsp)
            else:
                text = _collect_com_shape_text_for_markers(shape)
                if not text:
                    continue
                if "#CSP#" in text and deletecsp:
                    logging.debug("win32com: #CSP# 検出 → 削除 (shape): %r", text)
                if (
                    ("#CSP#" in text and deletecsp)
                    or "#TEMP#" in text
                    or "#MEMO#" in text
                ):
                    shape.Delete()
                    deleted += 1
        except Exception:
            pass
    return deleted


def modify_slide_shapes(pptx_path: Path, prs: Any, deletecsl: bool, deletecsp: bool) -> bool:
    # python-pptx pass: strip chapter/section/stay markers and do text cleanup.
    # Do NOT rely on python-pptx to decide whether win32com needs to run —
    # run-splitting and cross-paragraph text mean python-pptx may miss markers.
    has_temp_memo = False
    for slide in prs.slides:
        slide_full_text = "".join(
            frag for shape in slide.shapes for frag in _iter_shape_text_fragments(shape)
        )
        has_noreplace = NOREPLACE_MARKER in slide_full_text

        # Detect #TEMP#/#MEMO# via full slide text — covers text frames, tables, groups
        if "#TEMP#" in slide_full_text or "#MEMO#" in slide_full_text:
            has_temp_memo = True

        for shape in slide.shapes:
            if not has_noreplace:
                _strip_chapter_section_stay_markers_from_shape(shape)

            # --- table shapes ---
            if getattr(shape, "has_table", False) and shape.table is not None:
                if has_noreplace:
                    continue
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame is None:
                            continue
                        cell_text = cell.text_frame.text or ""
                        if "#TEMP#" in cell_text or "#MEMO#" in cell_text:
                            continue  # win32com will delete whole shape
                        if "#CSP#" in cell_text and deletecsp:
                            logging.debug("python-pptx: #CSP# 検出 (table cell): %r", cell_text)
                            continue  # win32com will delete whole shape
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                tobreak, _ = replace_CSLandCSP(run, deletecsl, deletecsp)
                                if tobreak:
                                    break
                continue

            # --- text frame shapes ---
            if not shape.has_text_frame:
                continue

            full_text = shape.text_frame.text or ""

            if "#TEMP#" in full_text or "#MEMO#" in full_text:
                continue  # skip text cleanup; win32com will delete this shape

            # deletecsp shape deletion is not skipped by NOREPLACE
            if "#CSP#" in full_text and deletecsp:
                logging.debug("python-pptx: #CSP# 検出 (shape): %r", full_text)
                continue  # will be deleted by win32com; skip text cleanup

            if has_noreplace:
                continue  # text cleanup skipped for NOREPLACE slides

            # Text cleanup (remove markers from shapes that are kept)
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    tobreak, _ = replace_CSLandCSP(run, deletecsl, deletecsp)
                    if tobreak:
                        break

    prs.save(str(pptx_path))

    # win32com pass: authoritative deletion using TextRange.Text (full concatenated text).
    # Run when any deletion flag is active or #TEMP#/#MEMO# shapes were found.
    # win32com itself decides which shapes to actually delete.
    needs_com = deletecsp or has_temp_memo
    if not needs_com:
        return False

    newpptx_path_abs = os.path.abspath(pptx_path)
    win32pptapp = win32com.client.Dispatch("PowerPoint.Application")
    win32prs = win32pptapp.Presentations.Open(str(newpptx_path_abs), WithWindow=False)
    deleted = 0
    for slide in win32prs.Slides:
        deleted += _delete_com_shapes_with_markers(slide.Shapes, deletecsp)
    win32prs.Save()
    win32prs.Close()
    logging.info("modify_slide_shapes: %d shape(s) deleted via win32com", deleted)
    return True


# ── idresolve モード ───────────────────────────────────────────────────────

_ID_MARKER_RE = re.compile(r'\{\{#id:([^}]+)\}\}')


def _parse_id_markers(notes_text: str) -> list[tuple[str, bool]]:
    """Parse {{#id:...}} markers from notes text line by line.

    Returns list of (full_id, nolabel) where nolabel is True when
    ``{{nolabel}}`` appears on the same line as the marker.
    """
    result = []
    for line in notes_text.splitlines():
        nolabel = "{{nolabel}}" in line
        for m in _ID_MARKER_RE.finditer(line):
            result.append((m.group(1), nolabel))
    return result


def _load_resolved_map(resolved_path: Path) -> dict[str, list[dict]]:
    with resolved_path.open(encoding="utf-8") as f:
        data = yaml.safe_load(f)
    result: dict[str, list[dict]] = {}
    for entry in (data.get("entries") or []):
        fid = entry.get("full_id", "")
        if fid:
            result.setdefault(fid, []).append(entry)
    return result


def _get_slide_notes_text(slide: Any) -> str:
    try:
        if not slide.has_notes_slide:
            return ""
        return slide.notes_slide.notes_text_frame.text or ""
    except Exception:
        return ""


def postprocess_idresolve(
    pptx_path: Path,
    resolved_path: Path,
    log_file: Path,
    duplicate_to: TextIO | None = None,
    config: dict[str, Any] | None = None,
) -> None:
    """Resolve slide titles from {{#id:...}} markers in notes using resolved ID map."""
    id_map = _load_resolved_map(resolved_path)
    prs = Presentation(str(pptx_path))

    lines: list[str] = []
    modified = False

    for slide_num, slide in enumerate(prs.slides, start=1):
        notes_text = _get_slide_notes_text(slide)
        markers = _parse_id_markers(notes_text)

        slide_full_text = "".join(
            frag for shape in slide.shapes for frag in _iter_shape_text_fragments(shape)
        )
        has_noreplace = NOREPLACE_MARKER in slide_full_text
        if has_noreplace and markers:
            logging.debug(
                "idresolve: スライド %d: %s があるため置換をスキップ",
                slide_num, NOREPLACE_MARKER,
            )

        for full_id, nolabel in markers:
            if has_noreplace:
                continue
            if nolabel:
                logging.debug(
                    "idresolve: スライド %d: {{nolabel}} があるためタイトル置換をスキップ (%s)",
                    slide_num, full_id,
                )
                continue
            matches = id_map.get(full_id, [])
            if len(matches) == 0:
                msg = f"[警告] スライド {slide_num}: full_id '{full_id}' が resolved に見つかりません"
                logging.warning("idresolve: slide %d: full_id '%s' not found in resolved", slide_num, full_id)
                lines.append(msg)
            elif len(matches) > 1:
                msg = f"[警告] スライド {slide_num}: full_id '{full_id}' が resolved に {len(matches)} 件あります"
                logging.warning("idresolve: slide %d: full_id '%s' has %d entries", slide_num, full_id, len(matches))
                lines.append(msg)
            else:
                label = matches[0].get("label", "")
                try:
                    slide.shapes.title.text = label
                    modified = True
                    lines.append(f"[{slide_num}] {label}")
                except Exception as exc:
                    msg = f"[警告] スライド {slide_num}: タイトル置換失敗 ({full_id}) - {exc}"
                    logging.warning("idresolve: slide %d: title replace failed: %s", slide_num, exc)
                    lines.append(msg)

    ppt_saved_by_com_shape_delete = _run_modify_slide_shapes(pptx_path, prs, config)

    if modified:
        if not ppt_saved_by_com_shape_delete:
            prs.save(str(pptx_path))

    block = "出力スライド一覧 (idresolve モード):\n" + "\n".join(lines) + "\n"
    if duplicate_to is not None:
        print(block, end="", file=duplicate_to)

    log_file.parent.mkdir(parents=True, exist_ok=True)
    with log_file.open("a", encoding="utf-8") as f:
        f.write(block)
