from __future__ import annotations

import argparse
import logging
import re
import shutil
import sys
from pathlib import Path
from typing import Any

import pythoncom
import time
import win32com.client
import yaml
from lxml import etree
from pptx import Presentation as _PptxPresentation
from pptx.oxml.ns import qn as _qn

_A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
_SOFT_RETURN_RE = re.compile(r'\x0B')
_INVALID_XML_RE = re.compile(r'[\x00-\x08\x0C\x0E-\x1F\x7F]')


def _sanitize_note_text(text: str) -> str:
    """Sanitize note text for XML: convert soft returns to newlines, strip other control chars."""
    text = _SOFT_RETURN_RE.sub('\n', text)
    return _INVALID_XML_RE.sub('', text)

from ._config import (
    ConfigError,
    load_yaml,
    validate_config,
    resolve_base,
    resolve_dir,

    resolve_file,
    resolve_file_in_dirs,
    parse_slides,
    parse_log_duplicate,
)
from md_merge.merge._recipe import _AUTO_FILENAMES, get_targetbasefilename, resolve_input_dirs
from ._slide_operation import SlideOperation, build_slide_insert_operations
from ._slide_titles import (
    _get_indexer_value,
    filter_slides_excluding_text_marker,
    NOREPLACE_MARKER,
)


def _apply_cli_indexer_opts(args: argparse.Namespace, recipe: dict[str, Any]) -> None:
    """Override recipe indexer flags with CLI options when explicitly passed."""
    if getattr(args, "delete_csl", False):
        recipe.setdefault("indexer", {})["deletecsl"] = True
    if getattr(args, "delete_csp", False):
        recipe.setdefault("indexer", {})["deletecsp"] = True
    pptxnumbering = getattr(args, "pptxnumbering", None)
    if pptxnumbering is not None:
        recipe.setdefault("indexer", {})["pptxnumbering"] = pptxnumbering


def format_slides_for_log(slide_list: list[int] | None) -> str:
    if slide_list is None:
        return "全スライド"
    if not slide_list:
        return "(挿入なし)"
    parts: list[str] = []
    for start, end in group_contiguous(slide_list):
        if start == end:
            parts.append(str(start))
        else:
            parts.append(f"{start}-{end}")
    return ", ".join(parts)


def resolve_merge_log_path(
    config_path: Path,
    config: dict[str, Any],
    output_dir: Path,
    output_filename: str,
    workdir: Path | None = None,
) -> Path:
    log_cfg = config.get("log") or {}
    filename = log_cfg.get("filename")
    if filename and str(filename).upper() != "STDOUT":
        log_dir_val = log_cfg.get("dir")
        if log_dir_val:
            log_dir = resolve_dir(log_dir_val, config_path, workdir)
        else:
            log_dir = output_dir
        log_dir.mkdir(parents=True, exist_ok=True)
        return log_dir / str(filename)
    output_dir.mkdir(parents=True, exist_ok=True)
    return output_dir / f"{Path(output_filename).stem}_merge.log"


def configure_merge_logging(
    config_path: Path,
    config: dict[str, Any],
    output_dir: Path,
    output_filename: str,
    workdir: Path | None = None,
) -> logging.Logger:
    log_cfg = config.get("log") or {}
    level_name = str(log_cfg.get("level", "info")).upper()
    level = getattr(logging, level_name, logging.INFO)

    log_path = resolve_merge_log_path(config_path, config, output_dir, output_filename, workdir)
    logger = logging.getLogger("pptx_merge.merge")
    logger.setLevel(level)
    logger.handlers.clear()
    logger.propagate = False

    fmt = logging.Formatter("%(asctime)s - %(levelname)s - %(message)s")
    handler = logging.FileHandler(log_path, encoding="utf-8")
    handler.setFormatter(fmt)
    logger.addHandler(handler)

    dup = parse_log_duplicate(config)
    if dup == "stdout":
        stream_h = logging.StreamHandler(sys.stdout)
        stream_h.setFormatter(fmt)
        logger.addHandler(stream_h)
    elif dup == "stderr":
        stream_h = logging.StreamHandler(sys.stderr)
        stream_h.setFormatter(fmt)
        logger.addHandler(stream_h)

    # Attach any general file handler from root logger so pptmerge internal
    # logs also appear in the recipe-level log file (log.filename).
    for h in logging.getLogger().handlers:
        if isinstance(h, logging.FileHandler):
            logger.addHandler(h)

    return logger


def group_contiguous(nums: list[int]) -> list[tuple[int, int]]:
    if not nums:
        return []
    groups = []
    start = prev = nums[0]
    for n in nums[1:]:
        if n == prev + 1:
            prev = n
        else:
            groups.append((start, prev))
            start = prev = n
    groups.append((start, prev))
    return groups


def create_destination_presentation(
    ppt_app: Any,
    config: dict[str, Any],
    config_path: Path,
    output_path: Path,
    workdir: Path | None = None,
    fallback_base: Path | None = None,
) -> Any:
    base = config.get("pptmerge") or {}
    base_file = base.get("stylebase")
    if base_file:
        base_dir = resolve_base(config_path, workdir)
        base_path = Path(base_file)
        if not base_path.is_absolute():
            base_path = base_dir / base_path
        if not base_path.exists():
            raise ConfigError(f"pptmerge.stylebase が見つかりません: {base_path}")
        shutil.copy2(base_path, output_path)
        return ppt_app.Presentations.Open(str(output_path), WithWindow=True)
    if fallback_base is not None:
        if not fallback_base.exists():
            raise ConfigError(f"スタイルベース（最初の insertpptx）が見つかりません: {fallback_base}")
        shutil.copy2(fallback_base, output_path)
        return ppt_app.Presentations.Open(str(output_path), WithWindow=True)
    return ppt_app.Presentations.Add()


def clear_slides(pres: Any) -> None:
    while pres.Slides.Count > 0:
        pres.Slides(1).Delete()


def add_blank_slide(dst: Any) -> None:
    pp_layout_blank = 12
    dst.Slides.Add(dst.Slides.Count + 1, pp_layout_blank)


def insert_slides_from_file(dst: Any, src_path: Path, slide_list: list[int] | None) -> None:
    insert_after = dst.Slides.Count
    if slide_list is None:
        dst.Slides.InsertFromFile(str(src_path), insert_after)
        return
    for start, end in group_contiguous(slide_list):
        insert_after = dst.Slides.Count
        dst.Slides.InsertFromFile(str(src_path), insert_after, start, end)


def _com_slide_title_text(slide: Any) -> str:
    """Return the slide title placeholder text, or \"\" if absent."""
    shapes = slide.Shapes
    if not shapes.HasTitle:
        return ""
    try:
        return str(shapes.Title.TextFrame.TextRange.Text or "")
    except Exception:
        return ""


def _collect_com_shape_text(shape: Any) -> str:
    """Return concatenated text from a COM shape (text frame, table, group children)."""
    parts: list[str] = []
    try:
        if shape.Type == 6:  # msoGroup
            for child in shape.GroupItems:
                parts.append(_collect_com_shape_text(child))
            return "".join(parts)
    except Exception:
        return ""
    try:
        if shape.HasTextFrame and shape.TextFrame.HasText:
            parts.append(shape.TextFrame.TextRange.Text or "")
    except Exception:
        pass
    try:
        if shape.HasTable:
            table = shape.Table
            for r in range(1, table.Rows.Count + 1):
                for c in range(1, table.Columns.Count + 1):
                    try:
                        parts.append(table.Cell(r, c).Shape.TextFrame.TextRange.Text or "")
                    except Exception:
                        pass
    except Exception:
        pass
    return "".join(parts)


def _replace_in_com_shape(shape: Any, find: str, value: str) -> None:
    """Replace find->value in all text of a COM shape (text frame, table, group)."""
    try:
        if shape.Type == 6:  # msoGroup
            for child in shape.GroupItems:
                _replace_in_com_shape(child, find, value)
            return
    except Exception:
        return
    try:
        if shape.HasTextFrame:
            shape.TextFrame.TextRange.Replace(find, value)
    except Exception:
        pass
    try:
        if shape.HasTable:
            table = shape.Table
            for r in range(1, table.Rows.Count + 1):
                for c in range(1, table.Columns.Count + 1):
                    try:
                        table.Cell(r, c).Shape.TextFrame.TextRange.Replace(find, value)
                    except Exception:
                        pass
    except Exception:
        pass


def _slide_has_noreplace_com(slide: Any) -> bool:
    """Return True if any shape on the slide contains NOREPLACE_MARKER."""
    for shape in slide.Shapes:
        if NOREPLACE_MARKER in _collect_com_shape_text(shape):
            return True
    return False


def _apply_notes_replacements(
    pptx_path: Path,
    pending: list[dict],
    resolved_map: dict[str, dict] | None = None,
) -> None:
    """Replace placeholders in slide notes via python-pptx (all occurrences, Python str.replace).

    Called after COM SaveAs so the file is closed.  Each entry in *pending* is a dict with:
      composedslidenum, basefilename, filename, slidenum, title, vars, skip
    """
    if not pending:
        return

    prs = _PptxPresentation(str(pptx_path))
    changed = False

    for info in pending:
        if info["skip"]:
            continue
        slide = prs.slides[info["composedslidenum"] - 1]
        if not slide.has_notes_slide:
            continue
        tf = slide.notes_slide.notes_text_frame
        text = _sanitize_note_text(tf.text or "")

        new_text = text
        new_text = new_text.replace("<__BASEFILENAME__>", info["basefilename"])
        new_text = new_text.replace("<__FILENAME__>", info["filename"])
        new_text = new_text.replace("<__SLIDENUM__>", info["slidenum"])
        new_text = new_text.replace("<__TITLE__>", info["title"])
        for name, value in info["vars"].items():
            new_text = new_text.replace("{{v:" + name + "}}", value)
        if resolved_map:
            new_text = _resolve_ref_in_text(new_text, resolved_map)

        if new_text == text:
            continue

        # Write back: rebuild paragraphs (one per line) from the replaced text
        txBody = tf._txBody
        for p in list(txBody.findall(_qn('a:p'))):
            txBody.remove(p)
        for line in new_text.split('\n'):
            p_elem = etree.SubElement(txBody, f'{{{_A_NS}}}p')
            if line:
                r_elem = etree.SubElement(p_elem, f'{{{_A_NS}}}r')
                t_elem = etree.SubElement(r_elem, f'{{{_A_NS}}}t')
                t_elem.text = line
        changed = True

    if changed:
        prs.save(str(pptx_path))


def replace_basefilename_in_inserted_slides(
    dst: Any,
    slide_start: int,
    slide_end: int,
    basefilename: str,
    filename: str,
) -> None:
    """Replace ``<__BASEFILENAME__>``, ``<__FILENAME__>``, ``<__SLIDENUM__>``, ``<__TITLE__>`` in slide shapes (not notes).

    Notes replacement is handled separately via :func:`_apply_notes_replacements`.
    """
    base_find = "<__BASEFILENAME__>"
    filename_find = "<__FILENAME__>"
    slidenum_find = "<__SLIDENUM__>"
    title_find = "<__TITLE__>"
    for i in range(slide_start, slide_end + 1):
        slide = dst.Slides(i)
        if _slide_has_noreplace_com(slide):
            continue
        slidenum_value = str(i - slide_start + 1)
        title_text = _com_slide_title_text(slide)
        for shape in slide.Shapes:
            _replace_in_com_shape(shape, base_find, basefilename)
            _replace_in_com_shape(shape, filename_find, filename)
            _replace_in_com_shape(shape, slidenum_find, slidenum_value)
            _replace_in_com_shape(shape, title_find, title_text)


_VAR_REF_RE = re.compile(r'\{\{v:([A-Za-z0-9_]+)\}\}')
_REF_RE = re.compile(r'\{\{(num|title|label):([^:}]+:[^:}]+:[^}]+)\}\}')


def _resolve_ref_in_text(text: str, resolved_map: dict[str, dict]) -> str:
    """Replace {{num/title/label:FULL_ID}} references using *resolved_map*."""
    def replace(m: re.Match) -> str:
        ref_type, full_id = m.group(1), m.group(2)
        entry = resolved_map.get(full_id)
        if entry is None:
            logging.warning(
                "pptmerge: 参照 {{%s:%s}} が resolved に見つかりません — そのまま残します",
                ref_type, full_id,
            )
            return m.group(0)
        if ref_type == "num":
            return entry.get("number", "")
        if ref_type == "title":
            return entry.get("title", "")
        # label
        return entry.get("label", "")

    return _REF_RE.sub(replace, text)


def _warn_undefined_vars(text: str, vars_map: dict[str, str], slide_num: int) -> None:
    for name in _VAR_REF_RE.findall(text):
        if name not in vars_map:
            logging.warning(
                "スライド %d: 未定義の変数参照 {{v:%s}} — vars: セクションに定義がありません",
                slide_num, name,
            )


def replace_vars_in_inserted_slides(
    dst: Any,
    slide_start: int,
    slide_end: int,
    vars_map: dict[str, str],
) -> None:
    """Replace ``{{v:NAME}}`` placeholders in slide shapes (not notes) using *vars_map*.

    Notes replacement is handled separately via :func:`_apply_notes_replacements`.
    """
    for i in range(slide_start, slide_end + 1):
        slide = dst.Slides(i)
        if _slide_has_noreplace_com(slide):
            continue
        for shape in slide.Shapes:
            _warn_undefined_vars(_collect_com_shape_text(shape), vars_map, i)
        if not vars_map:
            continue
        for name, value in vars_map.items():
            find = "{{v:" + name + "}}"
            for shape in slide.Shapes:
                _replace_in_com_shape(shape, find, value)


def replace_refs_in_inserted_slides(
    dst: Any,
    slide_start: int,
    slide_end: int,
    resolved_map: dict[str, dict],
) -> None:
    """Replace ``{{num/title/label:FULL_ID}}`` in slide shapes (not notes) using *resolved_map*.

    Notes replacement is handled separately via :func:`_apply_notes_replacements`.
    """
    if not resolved_map:
        return
    for i in range(slide_start, slide_end + 1):
        slide = dst.Slides(i)
        if _slide_has_noreplace_com(slide):
            continue
        # Collect full slide text to find which refs actually appear
        full_text = "".join(_collect_com_shape_text(shape) for shape in slide.Shapes)
        refs = _REF_RE.findall(full_text)  # list of (ref_type, full_id) tuples
        for ref_type, full_id in set(refs):
            entry = resolved_map.get(full_id)
            if entry is None:
                logging.warning(
                    "pptmerge: スライド %d: 参照 {{%s:%s}} が resolved に見つかりません",
                    i, ref_type, full_id,
                )
                continue
            if ref_type == "num":
                value = entry.get("number", "")
            elif ref_type == "title":
                value = entry.get("title", "")
            else:
                value = entry.get("label", "")
            find = "{{" + ref_type + ":" + full_id + "}}"
            for shape in slide.Shapes:
                _replace_in_com_shape(shape, find, value)


def replace_all_in_inserted_slides(
    dst: Any,
    slide_start: int,
    slide_end: int,
    basefilename: str,
    filename: str,
    vars_map: dict[str, str],
    resolved_map: dict[str, dict],
) -> None:
    """Replace all placeholders in slide shapes (not notes) in a single pass per slide.

    Combines the work of replace_basefilename_in_inserted_slides,
    replace_vars_in_inserted_slides, and replace_refs_in_inserted_slides to
    minimise COM round-trips:
      - Shape text is collected once per slide (NOREPLACE check + var warning +
        ref detection share the same collected string).
      - Only find strings actually present in the slide text are passed to
        TextRange.Replace.
      - Shapes are iterated once and all replacements are applied in that loop.
    """
    for i in range(slide_start, slide_end + 1):
        slide = dst.Slides(i)
        shapes = list(slide.Shapes)

        # Collect all shape text once
        all_text = "".join(_collect_com_shape_text(shape) for shape in shapes)

        if NOREPLACE_MARKER in all_text:
            continue

        # Warn about undefined vars
        for name in _VAR_REF_RE.findall(all_text):
            if name not in vars_map:
                logging.warning(
                    "スライド %d: 未定義の変数参照 {{v:%s}} — vars: セクションに定義がありません",
                    i, name,
                )

        # Build (find, value) list — skip entries not present in all_text
        slidenum_value = str(i - slide_start + 1)
        title_text = _com_slide_title_text(slide)

        replacements: list[tuple[str, str]] = []
        for find, value in (
            ("<__BASEFILENAME__>", basefilename),
            ("<__FILENAME__>", filename),
            ("<__SLIDENUM__>", slidenum_value),
            ("<__TITLE__>", title_text),
        ):
            if find in all_text:
                replacements.append((find, value))

        for name, value in vars_map.items():
            find = "{{v:" + name + "}}"
            if find in all_text:
                replacements.append((find, value))

        if resolved_map:
            for ref_type, full_id in set(_REF_RE.findall(all_text)):
                entry = resolved_map.get(full_id)
                if entry is None:
                    logging.warning(
                        "pptmerge: スライド %d: 参照 {{%s:%s}} が resolved に見つかりません",
                        i, ref_type, full_id,
                    )
                    continue
                if ref_type == "num":
                    value = entry.get("number", "")
                elif ref_type == "title":
                    value = entry.get("title", "")
                else:
                    value = entry.get("label", "")
                replacements.append(("{{" + ref_type + ":" + full_id + "}}", value))

        if not replacements:
            continue

        # Single pass: apply all replacements to each shape
        for shape in shapes:
            for find, value in replacements:
                _replace_in_com_shape(shape, find, value)


def merge_pptx(
    config_path: Path,
    workdir: Path | None = None,
    force: bool = False,
    args: argparse.Namespace | None = None,
    config: dict | None = None,
) -> tuple[Path, list[SlideOperation]]:
    if config is None:
        config = load_yaml(config_path)
        if args is not None:
            _apply_cli_indexer_opts(args, config)
    validate_config(config)

    output = config["output"]
    input_ = config.get("input") or {}
    separator = config.get("separator", {})
    vars_map: dict[str, str] = {str(k): str(v) for k, v in (config.get("vars") or {}).items()}

    outdir_val = output.get("outputdir")
    output_dir = (
        resolve_dir(outdir_val, config_path, workdir)
        if outdir_val
        else resolve_base(config_path, workdir)
    )

    # Load resolved ID map for {{num/title/label:FULL_ID}} substitution
    resolved_map: dict[str, dict] = {}
    idresolved_filename = output.get("idresolvedfilename")
    if not idresolved_filename:
        _base = get_targetbasefilename(output)
        if _base:
            idresolved_filename = _AUTO_FILENAMES["idresolvedfilename"].format(base=_base)
    if idresolved_filename:
        resolved_path_obj = output_dir / str(idresolved_filename)
        if resolved_path_obj.exists():
            with resolved_path_obj.open(encoding="utf-8") as _f:
                _data = yaml.safe_load(_f)
            resolved_map = {
                e["full_id"]: e
                for e in (_data.get("entries") or [])
                if e.get("full_id")
            }
            logging.info(
                "pptmerge: resolved map をロード: %d エントリ (%s)",
                len(resolved_map), resolved_path_obj,
            )
        else:
            logging.warning(
                "pptmerge: idresolvedfilename が見つかりません。参照置換をスキップします: %s",
                resolved_path_obj,
            )
    input_dirs = resolve_input_dirs(config, config_path, workdir)

    output_filename = output.get("pptxfilename")
    if not output_filename:
        _base = get_targetbasefilename(output)
        if _base:
            output_filename = f"{_base}.pptx"
        else:
            raise ConfigError("output.pptxfilename または output.targetbasefilename が必要です。")
    if Path(output_filename).name != output_filename:
        raise ConfigError("output.pptxfilename にはディレクトリを含めない単一ファイル名を指定してください。")

    output_path = output_dir / output_filename
    tmp_path = output_dir / ("._tmp_" + output_filename)

    if output_path.exists() and not force:
        raise ConfigError(
            f"出力ファイルが既に存在します（上書きするには --force を使用してください）: {output_path}"
        )

    output_dir.mkdir(parents=True, exist_ok=True)

    logger = configure_merge_logging(config_path, config, output_dir, output_filename, workdir)

    separator_enabled = bool(separator.get("enabled", False))

    _t_total = time.perf_counter()
    _t = time.perf_counter()
    ppt_app = win32com.client.Dispatch("PowerPoint.Application")
    ppt_app.Visible = True
    logger.info("[timing] COM Dispatch: %.3fs", time.perf_counter() - _t)

    # When pptmerge.stylebase is absent, use the first insertpptx file as style base
    first_src_path: Path | None = None
    if not (config.get("pptmerge") or {}).get("stylebase"):
        for item in config["procedure"]:
            if item.get("operation", "insertpptx") == "insertpptx":
                fname = item.get("pptxfilename")
                if fname:
                    first_src_path = resolve_file_in_dirs(fname, input_dirs)
                    break

    dst = None
    operations: list[SlideOperation] = []
    pending_notes: list[dict] = []
    rename_done = False

    try:
        _t = time.perf_counter()
        dst = create_destination_presentation(
            ppt_app, config, config_path, tmp_path, workdir, fallback_base=first_src_path
        )
        clear_slides(dst)
        logger.info("[timing] create+clear: %.3fs", time.perf_counter() - _t)

        for index, item in enumerate(config["procedure"]):
            operation = item.get("operation", "")

            if operation == "chapter":
                if item.get("chapter") is None:
                    logger.warning("chapter 操作に chapter キーがありません（index=%d）", index)
                logger.info(
                    "chapter 操作（pptx 挿入なし）: chapter=%r section=%r note=%r",
                    item.get("chapter"),
                    item.get("section"),
                    item.get("note"),
                )
                operations.append(SlideOperation(
                    action="chapter",
                    value=item.get("chapter"),
                    source="", slidenum=0, slidetitle="", composedslidenum=0,
                ))
                continue

            if operation == "section":
                if item.get("section") is None:
                    logger.warning("section 操作に section キーがありません（index=%d）", index)
                logger.info(
                    "section 操作（pptx 挿入なし）: chapter=%r section=%r note=%r",
                    item.get("chapter"),
                    item.get("section"),
                    item.get("note"),
                )
                operations.append(SlideOperation(
                    action="section",
                    value=item.get("section"),
                    source="", slidenum=0, slidetitle="", composedslidenum=0,
                ))
                continue

            if operation == "subsection":
                if item.get("subsection") is None:
                    logger.warning("subsection 操作に subsection キーがありません（index=%d）", index)
                logger.info("subsection 操作（pptx 挿入なし）: subsection=%r", item.get("subsection"))
                operations.append(SlideOperation(
                    action="subsection",
                    value=item.get("subsection"),
                    source="", slidenum=0, slidetitle="", composedslidenum=0,
                ))
                continue

            if operation == "beginstay":
                logger.info("beginstay 操作（pptx 挿入なし）")
                operations.append(SlideOperation(
                    action="beginstay",
                    value="", source="", slidenum=0, slidetitle="", composedslidenum=0,
                ))
                continue

            if operation == "endstay":
                logger.info("endstay 操作（pptx 挿入なし）")
                operations.append(SlideOperation(
                    action="endstay",
                    value="", source="", slidenum=0, slidetitle="", composedslidenum=0,
                ))
                continue

            if operation == "insertpptx":
                entry_file = item.get("pptxfilename")
                if entry_file is None:
                    raise ConfigError(
                        "operation が insertpptx のときは pptxfilename に入力 pptx を指定してください。"
                    )

                src_path = resolve_file_in_dirs(entry_file, input_dirs)
                if not src_path.exists():
                    raise ConfigError(f"入力ファイルが見つかりません: {src_path}")

                sep_before = item.get("separator_before", None)
                if sep_before is None:
                    sep_before = separator_enabled

                if index > 0 and sep_before:
                    composed_start = dst.Slides.Count + 1
                    add_blank_slide(dst)
                    operations.append(SlideOperation(
                        action="blank", value="", source="",
                        slidenum=0, slidetitle="", composedslidenum=composed_start,
                    ))

                _t_file = time.perf_counter()

                slide_list = parse_slides(item.get("slides", "all"))
                delete_csl_slides = bool(_get_indexer_value(config, "deletecsl"))
                if delete_csl_slides:
                    _t = time.perf_counter()
                    slide_list = filter_slides_excluding_text_marker(src_path, slide_list)
                    logger.info("[timing] %s filter_csl: %.3fs", src_path.name, time.perf_counter() - _t)

                composed_start = dst.Slides.Count + 1
                local_vars = {str(k): str(v) for k, v in (item.get("vars") or {}).items()}
                effective_vars = {**vars_map, **local_vars}
                _t = time.perf_counter()
                new_ops = build_slide_insert_operations(entry_file, src_path, slide_list, composed_start)
                logger.info("[timing] %s build_ops(pptx読込): %.3fs", src_path.name, time.perf_counter() - _t)
                if effective_vars:
                    for op in new_ops:
                        if op.slidetitle:
                            for name, value in effective_vars.items():
                                op.slidetitle = op.slidetitle.replace("{{v:" + name + "}}", value)
                operations.extend(new_ops)
                _t = time.perf_counter()
                insert_slides_from_file(dst, src_path, slide_list)
                logger.info("[timing] %s InsertFromFile: %.3fs", src_path.name, time.perf_counter() - _t)
                _t = time.perf_counter()
                replace_all_in_inserted_slides(
                    dst, composed_start, dst.Slides.Count,
                    src_path.stem, src_path.name, effective_vars, resolved_map,
                )
                logger.info("[timing] %s replace_all: %.3fs", src_path.name, time.perf_counter() - _t)
                # Collect per-slide note replacement info (applied after COM SaveAs)
                _t = time.perf_counter()
                for op in new_ops:
                    dst_slide = dst.Slides(op.composedslidenum)
                    pending_notes.append({
                        "composedslidenum": op.composedslidenum,
                        "basefilename": src_path.stem,
                        "filename": src_path.name,
                        "slidenum": str(op.slidenum),
                        "title": op.slidetitle,
                        "vars": dict(effective_vars),
                        "skip": _slide_has_noreplace_com(dst_slide),
                    })
                logger.info("[timing] %s pending_notes収集: %.3fs", src_path.name, time.perf_counter() - _t)
                logger.info(
                    "挿入: %s スライド: %s [ファイル合計 %.3fs]",
                    src_path,
                    format_slides_for_log(slide_list),
                    time.perf_counter() - _t_file,
                )

        _t = time.perf_counter()
        dst.SaveAs(str(tmp_path))
        dst.Close()
        dst = None
        logger.info("[timing] SaveAs+Close: %.3fs", time.perf_counter() - _t)

        _t = time.perf_counter()
        _apply_notes_replacements(tmp_path, pending_notes, resolved_map)
        logger.info("[timing] notes_replacements: %.3fs", time.perf_counter() - _t)

        logger.info("[timing] 合計: %.3fs", time.perf_counter() - _t_total)
        tmp_path.replace(output_path)
        rename_done = True

    finally:
        if dst is not None:
            try:
                dst.Close()
            except Exception:
                pass
        # try:     # powerpointを終了しないようにする
        #    ppt_app.Quit()
        # except Exception:
        #     pass
        if not rename_done and tmp_path.exists():
            try:
                tmp_path.unlink()
            except Exception:
                pass
        for h in logger.handlers[:]:
            try:
                h.flush()
                h.close()
            except Exception:
                pass
            logger.removeHandler(h)

    return output_path, operations
