from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation

from ._config import ConfigError


@dataclass
class SlideOperation:
    action: str
    value: str
    source: str
    slidenum: int
    slidetitle: str
    composedslidenum: int


def _slide_title_text(slide: object) -> str:
    if slide.shapes.title is not None:
        return (slide.shapes.title.text or "").strip()
    return ""


def build_slide_insert_operations(
    source: str,
    src_path: Path,
    slide_list: list[int] | None,
    composed_start: int,
    action: str = "insert",
) -> list[SlideOperation]:
    prs = Presentation(str(src_path))
    slides_list = list(prs.slides)
    n = len(slides_list)

    if slide_list is None:
        indices = list(range(1, n + 1))
    else:
        indices = slide_list
        for i in indices:
            if not isinstance(i, int) or i < 1 or i > n:
                raise ConfigError(
                    f"slides のスライド番号がファイルの範囲外です: {src_path} （指定 {i}、ファイルは {n} 枚）"
                )

    return [
        SlideOperation(
            source=source,
            slidenum=idx,
            action="insertslide",
            value="",
            slidetitle=_slide_title_text(slides_list[idx - 1]),
            composedslidenum=composed_start + i,
        )
        for i, idx in enumerate(indices)
    ]


def format_slide_operations_block(operations: list[SlideOperation]) -> str:
    if not operations:
        return ""
    lines: list[str] = ["挿入スライド操作（入力ファイル・元番号）:"]
    for op in operations:
        title_disp = op.slidetitle if op.slidetitle else "(タイトルなし)"
        if op.action == "insertslide":
            lines.append(
                f"  {op.action}: {op.source} 元スライド{op.slidenum} "
                f"結合後{op.composedslidenum} — {title_disp}"
            )
        else:
            lines.append(
                f"  {op.action}: {op.source} スライド{op.slidenum} — {title_disp}"
            )
    return "\n".join(lines) + "\n\n"
