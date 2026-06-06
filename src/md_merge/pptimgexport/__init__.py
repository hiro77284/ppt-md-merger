from __future__ import annotations

import argparse
import json
import logging
import math
from io import BytesIO
from pathlib import Path

import fitz  # pymupdf
from PIL import Image
from pptx import Presentation

from md_merge._filters import replace_slide_placeholders
from md_merge._note_parser import parse_note_image_directives
from md_merge._output import EXIT_BAD_INPUT, EXIT_FAILURE, EXIT_NOT_FOUND, EXIT_OK, emit, setup_logging
from md_merge.merge._recipe import load_yaml, resolve_input_path, resolve_workdir, setup_recipe_file_logging
from md_merge.pptmerge._config import resolve_base

_DEFAULT_DPI = 150
_MODE_SLIDE = "slide"
_MODE_NAMED_SHAPE = "named_shape"
_MODE_SLIDE_WITH_RECTPCT = "slide_with_rectpct"
_KNOWN_MODES = {_MODE_SLIDE, _MODE_NAMED_SHAPE, _MODE_SLIDE_WITH_RECTPCT}


def _normalize_cropsrc(cropsrc: dict | list) -> list[dict]:
    """cropsrc を常にリストとして返す（辞書は単一要素リストに変換）。"""
    if isinstance(cropsrc, list):
        return cropsrc
    if isinstance(cropsrc, dict):
        return [cropsrc]
    return []


def _round_pct(value: float) -> int:
    """四捨五入で整数パーセントに変換する。"""
    return math.floor(value + 0.5)


def _to_ltrb_int(
    shape_left: int, shape_top: int, shape_width: int, shape_height: int,
    slide_width: int, slide_height: int,
) -> tuple[int, int, int, int]:
    l = _round_pct(shape_left / slide_width * 100)
    t = _round_pct(shape_top / slide_height * 100)
    r = _round_pct((shape_left + shape_width) / slide_width * 100)
    b = _round_pct((shape_top + shape_height) / slide_height * 100)
    return l, t, r, b


def _to_ltrb_float(
    shape_left: int, shape_top: int, shape_width: int, shape_height: int,
    slide_width: int, slide_height: int,
) -> tuple[float, float, float, float]:
    l = shape_left / slide_width * 100
    t = shape_top / slide_height * 100
    r = (shape_left + shape_width) / slide_width * 100
    b = (shape_top + shape_height) / slide_height * 100
    return l, t, r, b



def _resolve_cropsrc_path(val: str, yaml_path: Path, workdir: Path | None) -> Path:
    base = resolve_base(yaml_path, workdir)
    p = Path(val)
    return p if p.is_absolute() else (base / p).resolve()


def _auto_pdf_path(pptx_val: str, output_section: dict, yaml_path: Path, workdir: Path | None) -> Path:
    """Auto-generate PDF path as output.outputdir/<pptx_stem>.pdf."""
    base = resolve_base(yaml_path, workdir)
    outdir_val = output_section.get("outputdir")
    if outdir_val:
        p = Path(str(outdir_val))
        out_dir = p if p.is_absolute() else (base / p).resolve()
    else:
        out_dir = base
    return out_dir / (Path(pptx_val).stem + ".pdf")


def _crop_pdf_page(
    pdf_doc: fitz.Document,
    page_idx: int,
    l: float, t: float, r: float, b: float,
    dpi: int,
) -> Image.Image:
    page = pdf_doc[page_idx]
    scale = dpi / 72
    pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale))
    img = Image.open(BytesIO(pix.tobytes("png")))
    w, h = img.size
    box = (round(w * l / 100), round(h * t / 100), round(w * r / 100), round(h * b / 100))
    return img.crop(box)


def _fmt_pct(v: float) -> str:
    return f"{v:.2f}" if v != int(v) else str(int(v))


def run(args: argparse.Namespace) -> int:
    setup_logging(args.log_level)

    yaml_path = resolve_input_path(args)
    if yaml_path is None:
        logging.error("No input file specified")
        return EXIT_BAD_INPUT
    if not yaml_path.exists():
        logging.error("Input file not found: %s", yaml_path)
        return EXIT_NOT_FOUND

    recipe = load_yaml(yaml_path)
    raw_cropsrc = (recipe.get("input") or {}).get("cropsrc") or {}
    entries = _normalize_cropsrc(raw_cropsrc)
    output_section = recipe.get("output") or {}
    cli_workdir = Path(args.workdir).resolve() if getattr(args, "workdir", None) else None
    workdir = resolve_workdir(recipe, yaml_path, cli_workdir)
    setup_recipe_file_logging(recipe, yaml_path, workdir)
    recipe_force = bool(output_section.get("force"))
    force = getattr(args, "force", False) or recipe_force
    dry_run = getattr(args, "dry_run", False)
    json_mode = getattr(args, "json", False)

    if not entries:
        logging.error("input.cropsrc が指定されていません: %s", yaml_path)
        return EXIT_BAD_INPUT

    # ── 全エントリのパス解決と存在確認（fail fast）────────────────────────
    resolved: list[tuple[Path, Path, int]] = []
    for i, entry in enumerate(entries):
        label = f"input.cropsrc[{i}]" if len(entries) > 1 else "input.cropsrc"
        pptx_val = entry.get("pptx")
        pdf_val = entry.get("pdf")
        dpi = int(entry.get("dpi", _DEFAULT_DPI))

        if not pptx_val:
            logging.error("%s.pptx が指定されていません: %s", label, yaml_path)
            return EXIT_BAD_INPUT
        pptx_path = _resolve_cropsrc_path(str(pptx_val), yaml_path, workdir)
        if pdf_val:
            pdf_path = _resolve_cropsrc_path(str(pdf_val), yaml_path, workdir)
        else:
            pdf_path = _auto_pdf_path(str(pptx_val), output_section, yaml_path, workdir)
            logging.debug("%s.pdf: 自動生成 %s", label, pdf_path)

        if not pptx_path.exists():
            logging.error("PPTX ファイルが見つかりません: %s", pptx_path)
            return EXIT_NOT_FOUND
        if not pdf_path.exists():
            logging.error("PDF ファイルが見つかりません: %s", pdf_path)
            return EXIT_NOT_FOUND

        resolved.append((pptx_path, pdf_path, dpi))

    # ── エントリごとに処理 ────────────────────────────────────────────────
    errors = 0
    json_entries: list[dict] = []

    for pptx_path, pdf_path, dpi in resolved:
        # PPTX 読み込み・タスク収集
        try:
            prs = Presentation(str(pptx_path))
        except Exception as e:
            logging.error("PPTX ファイルを開けません: %s: %s", pptx_path, e)
            errors += 1
            continue

        sw = prs.slide_width
        sh = prs.slide_height

        tasks: list[dict] = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            note_text = ""
            if slide.has_notes_slide:
                tf = slide.notes_slide.notes_text_frame
                note_text = (tf.text or "") if tf else ""

            title_shape = slide.shapes.title
            slide_title = title_shape.text if title_shape is not None else ""
            note_text = replace_slide_placeholders(note_text, pptx_path, slide_title, slide_idx)

            image_directives = parse_note_image_directives(note_text)
            for meta in image_directives:
                mode = meta.get("export-image-mode", "").strip()
                if mode not in _KNOWN_MODES:
                    continue

                target = meta.get("export-image-target", "").strip()
                export_val = meta.get("export-image", "").strip()

                if not export_val:
                    logging.warning("スライド %d: export-image が未指定です", slide_idx)
                    continue

                if mode == _MODE_SLIDE:
                    l: float = 0.0
                    t: float = 0.0
                    r: float = 100.0
                    b: float = 100.0
                    shape_label = _MODE_SLIDE
                elif mode == _MODE_NAMED_SHAPE:
                    if not target:
                        logging.warning("スライド %d: export-image-target が未指定です（mode: named_shape）", slide_idx)
                        continue
                    found = next((s for s in slide.shapes if s.name == target), None)
                    if found is None:
                        logging.warning("スライド %d: シェイプ '%s' が見つかりません", slide_idx, target)
                        continue
                    l, t, r, b = _to_ltrb_float(found.left, found.top, found.width, found.height, sw, sh)
                    shape_label = target
                else:  # slide_with_rectpct
                    if not target:
                        logging.warning("スライド %d: export-image-target が未指定です（mode: slide_with_rectpct）", slide_idx)
                        continue
                    found = next((s for s in slide.shapes if s.name == target), None)
                    if found is None:
                        logging.warning("スライド %d: シェイプ '%s' が見つかりません", slide_idx, target)
                        continue
                    li, ti, ri, bi = _to_ltrb_int(found.left, found.top, found.width, found.height, sw, sh)
                    l, t, r, b = float(li), float(ti), float(ri), float(bi)
                    shape_label = target

                export_path = Path(export_val)
                out_path = export_path if export_path.is_absolute() else (pptx_path.parent / export_path).resolve()

                tasks.append({
                    "slide": slide_idx,
                    "mode": mode,
                    "shape_label": shape_label,
                    "L": l, "T": t, "R": r, "B": b,
                    "output": out_path,
                })

        if not tasks:
            if not json_mode:
                print(f"対象スライドが見つかりません: {pptx_path}")
            json_entries.append({"pptx": str(pptx_path), "pdf": str(pdf_path), "tasks": [], "crops": []})
            continue

        # 座標一覧表示
        if not json_mode:
            print(f"対象スライド ({len(tasks)} 件): {pptx_path}")
            print(f"{'スライド':>6}  {'モード':<20}  {'シェイプ':<18}  {'L':>7}  {'T':>7}  {'R':>7}  {'B':>7}  出力先")
            print("-" * 110)
            for task in tasks:
                print(
                    f"{task['slide']:>6}  {task['mode']:<20}  {task['shape_label']:<18}  "
                    f"{_fmt_pct(task['L']):>7}  {_fmt_pct(task['T']):>7}  "
                    f"{_fmt_pct(task['R']):>7}  {_fmt_pct(task['B']):>7}  {task['output']}"
                )

        # PDF クロップ
        try:
            pdf_doc = fitz.open(str(pdf_path))
        except Exception as e:
            logging.error("PDF ファイルを開けません: %s: %s", pdf_path, e)
            errors += 1
            continue

        pdf_page_count = len(pdf_doc)
        crop_results: list[dict] = []

        for task in tasks:
            slide_n = task["slide"]
            page_idx = slide_n - 1
            out_path: Path = task["output"]

            if page_idx >= pdf_page_count:
                logging.warning(
                    "PDF にページ %d が存在しません（PDF は %d ページ）: %s",
                    slide_n, pdf_page_count, pdf_path,
                )
                errors += 1
                continue

            if out_path.exists() and not force:
                logging.error("出力ファイルが既に存在します (--force で上書き): %s", out_path)
                errors += 1
                continue

            if dry_run:
                emit(args, "ok", "pptimgexport", pdf_path, out_path, dry_run=True)
                crop_results.append({"slide": slide_n, "output": str(out_path)})
                continue

            try:
                out_path.parent.mkdir(parents=True, exist_ok=True)
                cropped = _crop_pdf_page(
                    pdf_doc, page_idx, task["L"], task["T"], task["R"], task["B"], dpi,
                )
                cropped.save(str(out_path))
                emit(args, "ok", "pptimgexport", pdf_path, out_path)
                crop_results.append({"slide": slide_n, "output": str(out_path)})
            except Exception as e:
                logging.error("crop 失敗: slide %d: %s", slide_n, e)
                errors += 1

        if getattr(args, "closepptx", False):
            pdf_doc.close()

        json_entries.append({
            "pptx": str(pptx_path),
            "pdf": str(pdf_path),
            "tasks": [_task_to_json(t) for t in tasks],
            "crops": crop_results,
        })

    if json_mode:
        print(json.dumps({
            "status": "ok" if errors == 0 else "error",
            "command": "pptimgexport",
            "entries": json_entries,
        }, ensure_ascii=False, indent=2))

    return EXIT_OK if errors == 0 else EXIT_FAILURE


def _task_to_json(task: dict) -> dict:
    return {**{k: v for k, v in task.items() if k != "output"}, "output": str(task["output"])}
